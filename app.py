import streamlit as st 
from helper import get_answer

from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
import streamlit as st

def file_extractor(files) : 
    '''
    Function to extract the text from the files

    Args:
    files : list : The list of files

    Returns:
    str : The text extracted from the files
    '''

    text = ''

    for fil in files : 

        if fil.name.endswith('doc') or fil.name.endswith('docx') : 

            for paragraph in Document(fil).paragraphs : 
                text += paragraph.text + '\n'

        elif fil.name.endswith('pdf') : 
                
                for page in PdfReader(fil).pages : 
                    text += page.extract_text() 
        elif fil.name.endswith('ppt') or fil.name.endswith('pptx') : 
                
                for slide in Presentation(fil).slides : 
                    for shape in slide.shapes : 
                        if hasattr(shape , 'text') : 
                            text += shape.text + '\n'

        return text

def process_input(question , text = open('sample.txt').read()):
    '''
    Function to process the input

    Args:
    question : str : The question to be asked
    text : str : The text from which the answer is to be extracted

    Returns:
    str : The answer to the question
    '''

    answer = get_answer(question , text)

    return answer


files = st.file_uploader('Uplaod your FIles' , accept_multiple_files = True , type = ['txt' , 'pdf' , 'docx' , 'pptx' , 'doc' , 'ppt']) 

if files : 

    text = file_extractor(files)

    question = st.text_input('Ask your question')

    if question : 
    
        answer = process_input(question , text)
        st.write(answer)
